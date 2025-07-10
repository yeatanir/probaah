# FILE: plugins/presentation/pptx_generator/research_slides.py
"""
Automated PowerPoint Generation for Research Presentations
Creates professional slides from analysis results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import json
from datetime import datetime
from typing import Dict, List, Optional
import pandas as pd

class ProbaahPresentationGenerator:
    """
    Automated research presentation generator
    Perfect for weekly updates, conferences, and group meetings
    """
    
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize presentation generator
        
        Args:
            template_path: Path to PowerPoint template (optional)
        """
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # Set up default styles
        self.setup_default_styles()
        
    def setup_default_styles(self):
        """Setup consistent styling for slides"""
        self.title_font_size = Pt(32)
        self.content_font_size = Pt(18)
        self.caption_font_size = Pt(14)
        
        # Penn State colors (or customize for your group)
        self.primary_color = RGBColor(30, 68, 128)  # Penn State Blue
        self.accent_color = RGBColor(180, 30, 50)   # Research Red
        self.text_color = RGBColor(50, 50, 50)      # Dark Gray
    
    def create_title_slide(self, title: str, subtitle: str = "", 
                          author: str = "Anirban Pal", 
                          affiliation: str = "van Duin Group, Penn State") -> None:
        """
        Create professional title slide
        
        Args:
            title: Presentation title
            subtitle: Subtitle or date
            author: Author name
            affiliation: Institution/group
        """
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        if slide.shapes.placeholders[1]:
            subtitle_shape = slide.shapes.placeholders[1]
            if subtitle:
                subtitle_text = f"{subtitle}\n\n{author}\n{affiliation}"
            else:
                subtitle_text = f"{author}\n{affiliation}\n{datetime.now().strftime('%B %d, %Y')}"
            
            subtitle_shape.text = subtitle_text
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = self.text_color
    
    def add_analysis_results_slide(self, results_file: str, 
                                  slide_title: str = "Analysis Results") -> None:
        """
        Create slide from analysis results JSON
        
        Args:
            results_file: Path to analysis results JSON
            slide_title: Title for the slide
        """
        with open(results_file, 'r') as f:
            results = json.load(f)
        
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Content
        content_shape = slide.shapes.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add key findings
        if 'bonds' in results:
            bonds = results['bonds']
            p = text_frame.paragraphs[0]
            p.text = f"• Average bond count: {bonds['avg_count']:.1f}"
            p.font.size = self.content_font_size
            
            p = text_frame.add_paragraph()
            p.text = f"• Average bond length: {bonds['avg_length']:.3f} Å"
            p.font.size = self.content_font_size
        
        if 'energy' in results and results['energy']:
            energy = results['energy']
            p = text_frame.add_paragraph()
            p.text = f"• Mean energy: {energy['mean']:.2f} eV"
            p.font.size = self.content_font_size
            
            p = text_frame.add_paragraph()
            p.text = f"• Energy stability: ±{energy['std']:.2f} eV"
            p.font.size = self.content_font_size
        
        if 'rdf' in results:
            rdf = results['rdf']
            p = text_frame.add_paragraph()
            p.text = f"• RDF analysis: {rdf['frames_analyzed']} frames analyzed"
            p.font.size = self.content_font_size
    
    def add_plot_slide(self, plot_file: str, title: str, 
                      caption: str = "", layout: str = "picture") -> None:
        """
        Add slide with plot/figure
        
        Args:
            plot_file: Path to plot image
            title: Slide title
            caption: Plot caption
            layout: Slide layout type
        """
        if layout == "picture":
            slide_layout = self.prs.slide_layouts[6]  # Blank slide
        else:
            slide_layout = self.prs.slide_layouts[1]  # Title and content
        
        slide = self.prs.slides.add_slide(slide_layout)
        
        if layout != "picture":
            # Add title
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Add image
        if Path(plot_file).exists():
            if layout == "picture":
                # Full slide image
                left = Inches(0.5)
                top = Inches(1.0)
                width = Inches(9)
                height = Inches(6.5)
            else:
                # Content area image
                left = Inches(1.5)
                top = Inches(2.5)
                width = Inches(7)
                height = Inches(4.5)
            
            pic = slide.shapes.add_picture(plot_file, left, top, width, height)
            
            # Add title on picture slides
            if layout == "picture":
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), 
                                                 Inches(9), Inches(0.8))
                text_frame = textbox.text_frame
                p = text_frame.paragraphs[0]
                p.text = title
                p.font.size = self.title_font_size
                p.font.color.rgb = self.primary_color
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            
            # Add caption
            if caption:
                if layout == "picture":
                    caption_top = Inches(7.8)
                else:
                    caption_top = Inches(7.2)
                
                textbox = slide.shapes.add_textbox(Inches(1), caption_top, 
                                                 Inches(8), Inches(0.5))
                text_frame = textbox.text_frame
                p = text_frame.paragraphs[0]
                p.text = caption
                p.font.size = self.caption_font_size
                p.font.color.rgb = self.text_color
                p.alignment = PP_ALIGN.CENTER
    
    def add_methods_slide(self, project_info: Dict) -> None:
        """
        Add methodology slide
        
        Args:
            project_info: Dictionary with project information
        """
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Methodology"
        title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Content
        content_shape = slide.shapes.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        methods = [
            "• Molecular Dynamics Simulations using ReaxFF",
            "• Force field: " + project_info.get('force_field', 'Custom ReaxFF parameters'),
            "• Temperature: " + project_info.get('temperature', '298 K'),
            "• Analysis performed using ASE and Probaah",
            "• Statistical analysis over " + str(project_info.get('n_frames', 'N')) + " frames"
        ]
        
        for i, method in enumerate(methods):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = method
            p.font.size = self.content_font_size
    
    def add_conclusions_slide(self, key_findings: List[str]) -> None:
        """
        Add conclusions slide
        
        Args:
            key_findings: List of key findings/conclusions
        """
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Key Findings & Conclusions"
        title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Content
        content_shape = slide.shapes.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, finding in enumerate(key_findings):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {finding}"
            p.font.size = self.content_font_size
            p.font.color.rgb = self.text_color
    
    def add_next_steps_slide(self, next_steps: List[str]) -> None:
        """
        Add next steps slide
        
        Args:
            next_steps: List of future work items
        """
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Next Steps"
        title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Content
        content_shape = slide.shapes.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, step in enumerate(next_steps):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {step}"
            p.font.size = self.content_font_size
            p.font.color.rgb = self.text_color
    
    def save_presentation(self, filename: str) -> str:
        """
        Save the presentation
        
        Args:
            filename: Output filename
            
        Returns:
            Path to saved presentation
        """
        output_path = Path(filename)
        self.prs.save(str(output_path))
        print(f"📊 Presentation saved: {output_path}")
        return str(output_path)

def create_weekly_update_presentation(analysis_dir: str, 
                                    output_file: str = "weekly_update.pptx",
                                    project_name: str = "Research Update") -> str:
    """
    Create automated weekly update presentation
    
    Args:
        analysis_dir: Directory containing analysis results
        output_file: Output PowerPoint file
        project_name: Name of the project
        
    Returns:
        Path to created presentation
    """
    generator = ProbaahPresentationGenerator()
    
    # Title slide
    date_str = datetime.now().strftime("%B %d, %Y")
    generator.create_title_slide(
        title=project_name,
        subtitle=f"Weekly Update - {date_str}",
        author="Anirban Pal",
        affiliation="van Duin Group, Penn State"
    )
    
    # Analysis results slide
    analysis_dir = Path(analysis_dir)
    results_file = analysis_dir / "analysis_results.json"
    if results_file.exists():
        generator.add_analysis_results_slide(str(results_file), "Analysis Results")
    
    # Plot slides
    plot_files = {
        "Bond Evolution": analysis_dir / "bond_evolution.png",
        "Radial Distribution Function": analysis_dir / "rdf.png", 
        "Energy Evolution": analysis_dir / "energy_evolution.png"
    }
    
    for title, plot_file in plot_files.items():
        if plot_file.exists():
            generator.add_plot_slide(str(plot_file), title, layout="picture")
    
    # Methods slide (if project config exists)
    config_file = analysis_dir.parent / ".probaah-config.yaml"
    if config_file.exists():
        import yaml
        with open(config_file, 'r') as f:
            config = yaml.safe_load(f)
        project_info = config.get('project', {})
        generator.add_methods_slide(project_info)
    
    # Conclusions (template - can be customized)
    key_findings = [
        "Successful molecular dynamics simulation completed",
        "Bond analysis reveals stable molecular structure",
        "Energy equilibration achieved within simulation timeframe",
        "Results consistent with experimental observations"
    ]
    generator.add_conclusions_slide(key_findings)
    
    # Next steps
    next_steps = [
        "Extend simulation time for longer-term analysis",
        "Investigate temperature effects on stability",
        "Compare with alternative force field parameters",
        "Prepare manuscript for publication"
    ]
    generator.add_next_steps_slide(next_steps)
    
    # Save presentation
    return generator.save_presentation(output_file)

# CLI interface function
def generate_presentation_cli(analysis_dir: str, title: str = "Research Presentation",
                            output: str = "presentation.pptx", 
                            style: str = "weekly") -> str:
    """
    Command-line interface for presentation generation
    
    Args:
        analysis_dir: Directory with analysis results
        title: Presentation title
        output: Output filename
        style: Presentation style (weekly, conference, group_meeting)
        
    Returns:
        Path to created presentation
    """
    if style == "weekly":
        return create_weekly_update_presentation(analysis_dir, output, title)
    else:
        # Can add other styles later
        return create_weekly_update_presentation(analysis_dir, output, title)

# Example usage
if __name__ == "__main__":
    # Test presentation generation
    # create_weekly_update_presentation("./analysis", "test_presentation.pptx")
    print("🎨 PowerPoint Generator ready for Probaah!")
